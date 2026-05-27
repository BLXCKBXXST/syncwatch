"""Build Sonar presentation from scratch.

Run:
    /tmp/pptx-venv/bin/python build.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

import style as s

ROOT = Path(__file__).parent
OUT_PPTX = ROOT / "final.pptx"

THEME = "Создание функциональной веб-платформы для потокового видео"
AUTHORS = [
    ("Андреев Андрей Валерьевич", "ст. преподаватель, СибГУТИ"),
    ("Лукинов Виталий Леонидович", "доцент, СибГУТИ"),
]


def set_bg(slide, color=s.BG_BASE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, x, y, w, h, text, font=s.FONT_SANS, size=s.BODY,
                bold=False, color=s.FG_PRIMARY, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = size
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_rect(slide, x, y, w, h, color):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    rect.shadow.inherit = False
    return rect


def add_header(slide, page_num):
    """Brand left, page number right, thin line below."""
    # left brand
    add_textbox(
        slide,
        s.PAD_X, s.HEADER_Y, Inches(9), Inches(0.3),
        s.BRAND_TEXT,
        font=s.FONT_MONO, size=s.HEADER_MONO, color=s.FG_MUTED,
    )
    # right page number
    add_textbox(
        slide,
        s.SLIDE_W - s.PAD_X - Inches(2), s.HEADER_Y, Inches(2), Inches(0.3),
        f"{page_num:02d} / {s.TOTAL_SLIDES:02d}",
        font=s.FONT_MONO, size=s.HEADER_MONO, color=s.FG_MUTED,
        align=PP_ALIGN.RIGHT,
    )
    # divider line
    line_w = s.SLIDE_W - 2 * s.PAD_X
    add_rect(slide, s.PAD_X, s.HEADER_LINE_Y, line_w, Emu(9525), s.BORDER)


def add_slide_title(slide, title):
    add_textbox(
        slide,
        s.PAD_X, s.TITLE_Y, s.CONTENT_W, s.TITLE_H,
        title,
        font=s.FONT_SANS, size=s.TITLE_H1, bold=True, color=s.ACCENT,
    )
    # short accent bar under title
    add_rect(
        slide,
        s.PAD_X, s.ACCENT_BAR_Y, s.ACCENT_BAR_W, s.ACCENT_BAR_H,
        s.ACCENT,
    )


# ---------- per-slide builders ----------

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    # tiny brand on top-left, no full header
    add_textbox(
        slide,
        s.PAD_X, s.HEADER_Y, Inches(6), Inches(0.3),
        "Sonar",
        font=s.FONT_MONO, size=Pt(13), color=s.ACCENT, bold=True,
    )
    # centered composition
    y0 = Inches(2.3)
    add_textbox(
        slide,
        s.PAD_X, y0, s.CONTENT_W, Inches(2.0),
        THEME,
        font=s.FONT_SANS, size=s.TITLE_HERO, bold=True,
        color=s.FG_PRIMARY, align=PP_ALIGN.LEFT,
    )
    # accent bar
    add_rect(slide, s.PAD_X, Inches(4.3), Inches(1.0), Inches(0.06), s.ACCENT)
    # authors
    y = Inches(4.6)
    for name, role in AUTHORS:
        add_textbox(
            slide,
            s.PAD_X, y, s.CONTENT_W, Inches(0.4),
            name,
            font=s.FONT_SANS, size=Pt(20), bold=True, color=s.FG_PRIMARY,
        )
        add_textbox(
            slide,
            s.PAD_X, y + Inches(0.35), s.CONTENT_W, Inches(0.35),
            role,
            font=s.FONT_SANS, size=s.META, color=s.FG_MUTED,
        )
        y += Inches(0.85)
    # date footer
    add_textbox(
        slide,
        s.PAD_X, s.SLIDE_H - Inches(0.6), s.CONTENT_W, Inches(0.3),
        "Итоговая работа · 2026",
        font=s.FONT_MONO, size=Pt(11), color=s.FG_MUTED,
    )


def slide_inner(prs, page_num, title):
    """Common inner-slide skeleton: header + title + accent bar."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, page_num)
    add_slide_title(slide, title)
    return slide


def slide_goal_tasks(prs, page_num):
    slide = slide_inner(prs, page_num, "Цель и задачи проекта")
    # goal block
    y = s.CONTENT_Y
    add_textbox(
        slide, s.PAD_X, y, s.CONTENT_W, Inches(0.35),
        "Цель проекта",
        font=s.FONT_MONO, size=Pt(12), color=s.FG_MUTED, bold=True,
    )
    add_textbox(
        slide, s.PAD_X, y + Inches(0.35), s.CONTENT_W, Inches(0.9),
        "Разработать функциональную веб-платформу для потокового видео\nс возможностью синхронного совместного просмотра.",
        font=s.FONT_SANS, size=s.BODY, color=s.FG_PRIMARY,
    )
    # tasks block
    y2 = y + Inches(1.5)
    add_textbox(
        slide, s.PAD_X, y2, s.CONTENT_W, Inches(0.35),
        "Задачи",
        font=s.FONT_MONO, size=Pt(12), color=s.FG_MUTED, bold=True,
    )
    tasks = [
        "1.  Спроектировать архитектуру системы: бэкенд, фронтенд, инфраструктура.",
        "2.  Реализовать потоковую отдачу видео по протоколу HLS.",
        "3.  Реализовать синхронные комнаты на WebSocket с серверной синхронизацией плеера.",
        "4.  Разработать адаптивный пользовательский интерфейс (десктоп и мобайл).",
        "5.  Обеспечить контейнерное развёртывание с автоматическим HTTPS.",
    ]
    add_textbox(
        slide, s.PAD_X, y2 + Inches(0.35), s.CONTENT_W, Inches(3.0),
        "\n".join(tasks),
        font=s.FONT_SANS, size=s.BODY, color=s.FG_SECONDARY,
    )


def slide_roles(prs, page_num):
    slide = slide_inner(prs, page_num, "Распределение ролей")
    rows = [
        ("Лукинов В.Л.",   "Моделирование архитектуры информационной системы"),
        ("Лукинов В.Л.",   "Разработка бизнес-логики информационной системы"),
        ("Андреев А.В.",   "Прототипирование и разработка пользовательского интерфейса"),
        ("Андреев А.В.",   "Тестирование и интеграционные работы"),
    ]
    # custom-built table via rectangles for full control
    col1_w = Inches(3.2)
    col2_w = s.CONTENT_W - col1_w
    header_h = Inches(0.55)
    row_h = Inches(0.85)

    y = s.CONTENT_Y
    # header
    add_rect(slide, s.PAD_X, y, s.CONTENT_W, header_h, s.ACCENT)
    add_textbox(
        slide, s.PAD_X + Inches(0.25), y, col1_w, header_h,
        "Фамилия И.О.",
        font=s.FONT_SANS, size=Pt(15), bold=True, color=s.WHITE,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        slide, s.PAD_X + col1_w + Inches(0.25), y, col2_w, header_h,
        "Роль в группе",
        font=s.FONT_SANS, size=Pt(15), bold=True, color=s.WHITE,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    y += header_h
    # rows
    for i, (name, role) in enumerate(rows):
        bg_color = s.BG_BASE if i % 2 == 0 else s.BG_MUTED
        add_rect(slide, s.PAD_X, y, s.CONTENT_W, row_h, bg_color)
        # bottom border
        add_rect(slide, s.PAD_X, y + row_h - Emu(6350), s.CONTENT_W, Emu(6350), s.BORDER)
        add_textbox(
            slide, s.PAD_X + Inches(0.25), y, col1_w, row_h,
            name,
            font=s.FONT_SANS, size=Pt(16), color=s.FG_PRIMARY,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_textbox(
            slide, s.PAD_X + col1_w + Inches(0.25), y, col2_w, row_h,
            role,
            font=s.FONT_SANS, size=Pt(15), color=s.FG_SECONDARY,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        y += row_h


def slide_means(prs, page_num):
    slide = slide_inner(prs, page_num, "Средства реализации")
    # two-column layout
    col_w = (s.CONTENT_W - Inches(0.5)) / 2
    col1 = [
        ("Бэкенд",            "Python 3.10+, Django 5.1, Django REST Framework 3.15, Django Channels 4.1 (ASGI + WebSocket)"),
        ("База данных",       "PostgreSQL 16 (продакшен), SQLite (разработка)"),
        ("Кэш и pub/sub",     "Redis 7 (channel layer для WebSocket)"),
        ("Потоковое видео",   "HLS, hls.js, HTML5 video"),
    ]
    col2 = [
        ("Фронтенд",          "React 19.2, React Router 7, Vite 8, axios"),
        ("Инфраструктура",    "Docker, docker-compose, Caddy (reverse proxy и auto-HTTPS Let's Encrypt)"),
        ("Тестирование",      "pytest + pytest-asyncio (бэк), Vitest + @testing-library/react (фронт)"),
        ("Аутентификация",    "JWT (SimpleJWT), вход по одноразовому коду на e-mail"),
    ]

    def render_col(x, items):
        y = s.CONTENT_Y
        for label, val in items:
            add_textbox(
                slide, x, y, col_w, Inches(0.3),
                label,
                font=s.FONT_MONO, size=Pt(11), color=s.ACCENT, bold=True,
            )
            add_textbox(
                slide, x, y + Inches(0.3), col_w, Inches(0.8),
                val,
                font=s.FONT_SANS, size=Pt(14), color=s.FG_PRIMARY,
            )
            y += Inches(1.25)

    render_col(s.PAD_X, col1)
    render_col(s.PAD_X + col_w + Inches(0.5), col2)


def slide_image(prs, page_num, title, caption=None, png_path=None, hint=None):
    """Слайд с заголовком и картинкой по центру.
    Если caption задан — подпись над картинкой, бокс 12.13" × 4.5" под акцент-баром.
    Если caption=None — без акцент-бара; зона картинки 12.73" × 5.6", левый паддинг 0.3"."""
    slide = slide_inner(prs, page_num, title)
    if caption:
        add_textbox(
            slide, s.PAD_X, s.CONTENT_Y, s.CONTENT_W, Inches(0.4),
            caption,
            font=s.FONT_SANS, size=s.BODY_SMALL, color=s.FG_MUTED,
        )
        box_x = s.PAD_X
        box_w = s.CONTENT_W
        box_y = s.CONTENT_Y + Inches(0.65)
        box_h = Inches(4.5)
    else:
        box_x = Inches(0.3)
        box_w = s.SLIDE_W - Inches(0.6)
        box_y = Inches(1.65)
        box_h = Inches(5.6)
    if png_path and Path(png_path).exists():
        from PIL import Image as _PILImage
        with _PILImage.open(png_path) as im:
            iw, ih = im.size
        # fit image into box preserving aspect ratio
        ratio = min(box_w / (iw * 9525), box_h / (ih * 9525))
        w = int(iw * 9525 * ratio)
        h = int(ih * 9525 * ratio)
        x = box_x + (box_w - w) // 2
        y = box_y + (box_h - h) // 2
        slide.shapes.add_picture(str(png_path), x, y, width=w, height=h)
    else:
        add_rect(slide, box_x, box_y, box_w, box_h, s.BG_ELEVATED)
        add_textbox(
            slide, box_x, box_y, box_w, box_h,
            hint or "[ контент ]",
            font=s.FONT_MONO, size=Pt(16), color=s.FG_MUTED,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )


def _draw_browser_frame(slide, x, y, w, h, url):
    """Browser-окно: внешняя рамка + chrome bar (3 точки + URL). Возвращает chrome_h."""
    add_rect(slide, x, y, w, h, s.BG_BASE)
    add_rect(slide, x, y, w, Emu(9525), s.BORDER_STRONG)
    add_rect(slide, x, y + h - Emu(9525), w, Emu(9525), s.BORDER_STRONG)
    add_rect(slide, x, y, Emu(9525), h, s.BORDER_STRONG)
    add_rect(slide, x + w - Emu(9525), y, Emu(9525), h, s.BORDER_STRONG)
    chrome_h = Inches(0.4)
    add_rect(slide, x, y, w, chrome_h, s.BG_MUTED)
    for j, color in enumerate([RGBColor(0xEF, 0x44, 0x44), RGBColor(0xF5, 0x9E, 0x0B), RGBColor(0x10, 0xB9, 0x81)]):
        dot_x = x + Inches(0.15) + j * Inches(0.22)
        dot_y = y + Inches(0.12)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, dot_x, dot_y, Inches(0.16), Inches(0.16))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        dot.shadow.inherit = False
    add_textbox(
        slide, x + Inches(1.0), y + Inches(0.05), w - Inches(1.2), chrome_h - Inches(0.05),
        url,
        font=s.FONT_MONO, size=Pt(12), color=s.FG_MUTED, anchor=MSO_ANCHOR.MIDDLE,
    )
    return chrome_h


def _draw_png_in_frame(slide, x, y, w, h, png_path, hint):
    """Вставляет PNG в зону x..w / y..h с сохранением пропорций. При отсутствии PNG — серый placeholder."""
    if png_path and Path(png_path).exists():
        from PIL import Image as _PILImage
        with _PILImage.open(png_path) as im:
            iw, ih = im.size
        ratio = min(w / (iw * 9525), h / (ih * 9525))
        rw = int(iw * 9525 * ratio)
        rh = int(ih * 9525 * ratio)
        rx = x + (w - rw) // 2
        ry = y + (h - rh) // 2
        slide.shapes.add_picture(str(png_path), rx, ry, width=rw, height=rh)
    else:
        add_rect(slide, x, y, w, h, s.BG_ELEVATED)
        add_textbox(
            slide, x, y, w, h, hint,
            font=s.FONT_MONO, size=Pt(11), color=s.FG_MUTED,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )


def _draw_sonar_header(slide, x, y, w):
    """Синяя полоса-шапка с «Sonar» внутри окна."""
    h = Inches(0.55)
    add_rect(slide, x, y, w, h, s.ACCENT)
    add_textbox(
        slide, x + Inches(0.3), y, w - Inches(0.6), h,
        "Sonar",
        font=s.FONT_MONO, size=Pt(20), bold=True, color=s.WHITE,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    return h


def _draw_mockup_home(slide, x, y, w, h):
    """Макет «Главная»: шапка Sonar + 2×3 сетка карточек видео с двумя короткими линиями-подписями."""
    head_h = _draw_sonar_header(slide, x, y, w)
    pad = Inches(0.4)
    cy_start = y + head_h + Inches(0.4)
    cols, rows = 3, 2
    gap_x, gap_y = Inches(0.3), Inches(0.45)
    card_w = (w - 2 * pad - (cols - 1) * gap_x) / cols
    card_h = int(card_w / 1.5)  # landscape aspect 1.5
    line_block_h = Inches(0.22)  # высота под две линии-подписи
    for row in range(rows):
        for col in range(cols):
            cx = x + pad + col * (card_w + gap_x)
            cy = cy_start + row * (card_h + line_block_h + gap_y)
            # карточка
            add_rect(slide, cx, cy, card_w, card_h, s.BG_ELEVATED)
            # две короткие линии одинаковой длины и толщины (без контраста)
            line_w = int(card_w * 0.6)
            add_rect(slide, cx, cy + card_h + Inches(0.08), line_w, Inches(0.06), s.BORDER_STRONG)
            add_rect(slide, cx, cy + card_h + Inches(0.18), int(card_w * 0.45), Inches(0.06), s.BORDER_STRONG)


def _draw_mockup_room(slide, x, y, w, h):
    """Макет «Комната просмотра»: шапка + 16:9 плеер + чат-сайдбар из узких карточек + поле ввода."""
    head_h = _draw_sonar_header(slide, x, y, w)
    pad = Inches(0.35)
    inner_x = x + pad
    inner_w = w - 2 * pad
    inner_y = y + head_h + Inches(0.4)

    # плеер 16:9 (горизонтальный прямоугольник, ~65% ширины)
    gap = Inches(0.2)
    player_w = int(inner_w * 0.65)
    player_h = int(player_w * 9 / 16)  # 16:9 aspect
    chat_w = inner_w - player_w - gap
    chat_h = player_h  # одинаковая высота с плеером

    # плеер
    add_rect(slide, inner_x, inner_y, player_w, player_h, s.FG_PRIMARY)
    # play-треугольник по центру
    tri_size = Inches(0.6)
    tri_x = inner_x + (player_w - tri_size) // 2
    tri_y = inner_y + (player_h - tri_size) // 2
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, tri_x, tri_y, tri_size, tri_size)
    tri.rotation = 90
    tri.fill.solid()
    tri.fill.fore_color.rgb = s.WHITE
    tri.line.fill.background()
    tri.shadow.inherit = False
    # прогресс-бар: толще, точно под шириной плеера
    add_rect(slide, inner_x, inner_y + player_h + Inches(0.08), player_w, Inches(0.1), s.ACCENT)

    # чат-сайдбар: 4 УЗКИЕ горизонтальные карточки в стек
    chat_x = inner_x + player_w + gap
    msg_gap = Inches(0.1)
    msg_h = (chat_h - 3 * msg_gap) / 4
    for k in range(4):
        add_rect(slide, chat_x, inner_y + k * (msg_h + msg_gap), chat_w, msg_h, s.BG_MUTED)

    # поле ввода — горизонтальная полоса на всю ширину inner_w
    input_y = inner_y + player_h + Inches(0.35)
    add_rect(slide, inner_x, input_y, inner_w, Inches(0.35), s.BG_ELEVATED)


def _draw_mockup_profile(slide, x, y, w, h):
    """Макет «Профиль»: шапка + плоский круг-аватар + длинная и короткая строки имени + тонкая tabs + 2×3 сетка."""
    head_h = _draw_sonar_header(slide, x, y, w)
    pad = Inches(0.4)
    cy = y + head_h + Inches(0.4)

    # аватар — плоский цвет без обводки
    av_size = Inches(0.9)
    av = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + pad, cy, av_size, av_size)
    av.fill.solid()
    av.fill.fore_color.rgb = s.ACCENT_SOFT
    av.line.fill.background()  # без обводки
    av.shadow.inherit = False
    # строки имени: тёмная длинная + серая короче
    name_x = x + pad + av_size + Inches(0.3)
    add_rect(slide, name_x, cy + Inches(0.15), Inches(2.4), Inches(0.2), s.FG_PRIMARY)
    add_rect(slide, name_x, cy + Inches(0.5), Inches(1.6), Inches(0.13), s.FG_MUTED)
    # tabs-индикатор: уже и тоньше
    tabs_y = cy + av_size + Inches(0.3)
    add_rect(slide, x + pad, tabs_y, Inches(0.55), Inches(0.06), s.ACCENT)

    # сетка карточек 3×2 с увеличенным воздухом
    grid_y = tabs_y + Inches(0.35)
    cols, rows = 3, 2
    gap_x, gap_y = Inches(0.3), Inches(0.35)
    card_w = (w - 2 * pad - (cols - 1) * gap_x) / cols
    grid_h_avail = y + h - grid_y - Inches(0.25)
    card_h = (grid_h_avail - (rows - 1) * gap_y) / rows
    for row in range(rows):
        for col in range(cols):
            cx = x + pad + col * (card_w + gap_x)
            cy2 = grid_y + row * (card_h + gap_y)
            add_rect(slide, cx, cy2, card_w, card_h, s.BG_ELEVATED)


def _slide_ui_pair(prs, page_num, title, caption, mock_renderer, real_png, mock_url, real_url):
    """Парный слайд: слева программный Figma-макет, справа реальный скрин платформы."""
    slide = slide_inner(prs, page_num, title)
    add_textbox(
        slide, s.PAD_X, s.CONTENT_Y, s.CONTENT_W, Inches(0.4),
        caption,
        font=s.FONT_SANS, size=s.BODY_SMALL, color=s.FG_MUTED,
    )
    gap = Inches(0.35)
    pad_x = Inches(0.3)
    mock_w = Inches(4.2)
    win_h = Inches(4.6)
    real_w = s.SLIDE_W - 2 * pad_x - gap - mock_w
    base_y = s.CONTENT_Y + Inches(0.55)
    left_x = pad_x
    right_x = left_x + mock_w + gap

    img_dir = ROOT / "img"

    # left: программный макет (почти квадратное окно)
    ch_l = _draw_browser_frame(slide, left_x, base_y, mock_w, win_h, mock_url)
    inner_l_y = base_y + ch_l
    inner_l_h = win_h - ch_l - Emu(9525)
    inner_l_x = left_x + Emu(9525)
    inner_l_w = mock_w - 2 * Emu(9525)
    mock_renderer(slide, inner_l_x, inner_l_y, inner_l_w, inner_l_h)

    # right: реальный скрин (широкое окно для детализации)
    ch_r = _draw_browser_frame(slide, right_x, base_y, real_w, win_h, real_url)
    inner_r_y = base_y + ch_r
    inner_r_h = win_h - ch_r - Emu(9525)
    inner_r_x = right_x + Emu(9525)
    inner_r_w = real_w - 2 * Emu(9525)
    _draw_png_in_frame(slide, inner_r_x, inner_r_y, inner_r_w, inner_r_h, img_dir / real_png, f"[ скриншот: {title.lower()} ]")

    cap_y = base_y + win_h + Inches(0.1)
    add_textbox(
        slide, left_x, cap_y, mock_w, Inches(0.3),
        "Макет (Figma)",
        font=s.FONT_MONO, size=Pt(12), color=s.FG_MUTED, align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, right_x, cap_y, real_w, Inches(0.3),
        "Реализация",
        font=s.FONT_MONO, size=Pt(12), color=s.FG_MUTED, align=PP_ALIGN.CENTER,
    )


def slide_ui_home(prs, page_num):
    _slide_ui_pair(
        prs, page_num,
        "Главная: макет и реализация",
        "Лента видео — главный экран платформы.",
        _draw_mockup_home, "ui-home.png",
        "figma.com/Sonar/home", "sonar.app/",
    )


def slide_ui_room(prs, page_num):
    _slide_ui_pair(
        prs, page_num,
        "Совместный просмотр: макет и реализация",
        "Комната с общим плеером, чатом и панелью вопросов.",
        _draw_mockup_room, "ui-room.png",
        "figma.com/Sonar/room", "sonar.app/room/<id>",
    )


def slide_ui_profile(prs, page_num):
    _slide_ui_pair(
        prs, page_num,
        "Профиль: макет и реализация",
        "Канал пользователя со списком загруженных видео.",
        _draw_mockup_profile, "ui-profile.png",
        "figma.com/Sonar/profile", "sonar.app/profile",
    )


def _slide_ui_mockup_LEGACY(prs, page_num):
    """Слайд 9: стилизованные плейсхолдеры под три экрана UI до готовности темы Sonar."""
    slide = slide_inner(prs, page_num, "Прототипирование пользовательского интерфейса")
    add_textbox(
        slide, s.PAD_X, s.CONTENT_Y, s.CONTENT_W, Inches(0.4),
        "Ключевые экраны платформы: главная, комната совместного просмотра, профиль.",
        font=s.FONT_SANS, size=s.BODY_SMALL, color=s.FG_MUTED,
    )
    # three browser-like windows
    gap = Inches(0.3)
    win_w = (s.CONTENT_W - 2 * gap) / 3
    win_h = Inches(3.6)
    base_y = s.CONTENT_Y + Inches(0.65)
    titles = ["Главная", "Комната просмотра", "Профиль"]
    for i, t in enumerate(titles):
        x = s.PAD_X + i * (win_w + gap)
        # outer border
        add_rect(slide, x, base_y, win_w, win_h, s.BG_BASE)
        add_rect(slide, x, base_y, win_w, Emu(9525), s.BORDER_STRONG)  # top
        add_rect(slide, x, base_y + win_h - Emu(9525), win_w, Emu(9525), s.BORDER_STRONG)
        add_rect(slide, x, base_y, Emu(9525), win_h, s.BORDER_STRONG)
        add_rect(slide, x + win_w - Emu(9525), base_y, Emu(9525), win_h, s.BORDER_STRONG)
        # mac-style chrome bar
        chrome_h = Inches(0.32)
        add_rect(slide, x, base_y, win_w, chrome_h, s.BG_MUTED)
        # three dots
        for j, color in enumerate([RGBColor(0xEF,0x44,0x44), RGBColor(0xF5,0x9E,0x0B), RGBColor(0x10,0xB9,0x81)]):
            dot_x = x + Inches(0.12) + j * Inches(0.18)
            dot_y = base_y + Inches(0.1)
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, dot_x, dot_y, Inches(0.13), Inches(0.13))
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.fill.background()
            dot.shadow.inherit = False
        # url-like text
        add_textbox(
            slide, x + Inches(0.9), base_y + Inches(0.04), win_w - Inches(1.0), chrome_h - Inches(0.04),
            "sonar.app/" + ("home" if i == 0 else "room/r-42" if i == 1 else "u/host"),
            font=s.FONT_MONO, size=Pt(9), color=s.FG_MUTED, anchor=MSO_ANCHOR.MIDDLE,
        )
        # inner header strip with Sonar
        inner_y = base_y + chrome_h
        add_rect(slide, x, inner_y, win_w, Inches(0.4), s.ACCENT)
        add_textbox(
            slide, x + Inches(0.18), inner_y, win_w - Inches(0.36), Inches(0.4),
            "Sonar",
            font=s.FONT_MONO, size=Pt(11), bold=True, color=s.WHITE,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # content skeleton — varying per screen
        cy = inner_y + Inches(0.55)
        if i == 0:  # Главная — карточки видео
            for row in range(2):
                for col in range(3):
                    cx = x + Inches(0.2) + col * Inches(0.95)
                    cy2 = cy + row * Inches(1.05)
                    add_rect(slide, cx, cy2, Inches(0.85), Inches(0.55), s.BG_ELEVATED)
                    add_rect(slide, cx, cy2 + Inches(0.6), Inches(0.7), Inches(0.06), s.BORDER_STRONG)
                    add_rect(slide, cx, cy2 + Inches(0.72), Inches(0.5), Inches(0.06), s.BORDER)
        elif i == 1:  # Комната — плеер + чат
            # video area
            player_w = win_w - Inches(1.6)
            add_rect(slide, x + Inches(0.2), cy, player_w, Inches(1.6), s.FG_PRIMARY)
            # play triangle
            tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, x + Inches(0.2) + player_w/2 - Inches(0.18), cy + Inches(0.55), Inches(0.36), Inches(0.5))
            tri.rotation = 90
            tri.fill.solid()
            tri.fill.fore_color.rgb = s.WHITE
            tri.line.fill.background()
            tri.shadow.inherit = False
            # progress bar
            add_rect(slide, x + Inches(0.2), cy + Inches(1.66), player_w, Inches(0.05), s.ACCENT)
            # chat sidebar
            chat_x = x + Inches(0.2) + player_w + Inches(0.1)
            for k in range(4):
                add_rect(slide, chat_x, cy + k * Inches(0.34), Inches(1.2), Inches(0.28), s.BG_MUTED)
            # input
            add_rect(slide, x + Inches(0.2), cy + Inches(1.9), win_w - Inches(0.4), Inches(0.3), s.BG_ELEVATED)
        else:  # Профиль — шапка + видео-сетка
            # avatar
            av = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), cy, Inches(0.55), Inches(0.55))
            av.fill.solid()
            av.fill.fore_color.rgb = s.ACCENT_SOFT
            av.line.color.rgb = s.ACCENT
            av.shadow.inherit = False
            add_rect(slide, x + Inches(0.85), cy + Inches(0.06), Inches(1.4), Inches(0.12), s.FG_PRIMARY)
            add_rect(slide, x + Inches(0.85), cy + Inches(0.24), Inches(1.0), Inches(0.08), s.FG_MUTED)
            # tabs
            add_rect(slide, x + Inches(0.2), cy + Inches(0.75), Inches(0.6), Inches(0.05), s.ACCENT)
            # grid
            for row in range(2):
                for col in range(3):
                    cx = x + Inches(0.2) + col * Inches(0.95)
                    cy2 = cy + Inches(0.95) + row * Inches(0.7)
                    add_rect(slide, cx, cy2, Inches(0.85), Inches(0.55), s.BG_ELEVATED)
        # caption под окном
        add_textbox(
            slide, x, base_y + win_h + Inches(0.1), win_w, Inches(0.35),
            t,
            font=s.FONT_MONO, size=Pt(12), color=s.FG_MUTED,
            align=PP_ALIGN.CENTER,
        )
    # footnote
    add_textbox(
        slide, s.PAD_X, s.SLIDE_H - Inches(0.45), s.CONTENT_W, Inches(0.3),
        "Макеты интерфейса в Figma.",
        font=s.FONT_MONO, size=Pt(10), color=s.FG_MUTED, align=PP_ALIGN.CENTER,
    )


def _slide_real_ui_LEGACY(prs, page_num):
    """Слайд 10: реальные скриншоты живой платформы в тех же browser-окнах."""
    slide = slide_inner(prs, page_num, "Реализация пользовательского интерфейса")
    add_textbox(
        slide, s.PAD_X, s.CONTENT_Y, s.CONTENT_W, Inches(0.4),
        "Реализованные экраны платформы: главная, комната совместного просмотра, профиль.",
        font=s.FONT_SANS, size=s.BODY_SMALL, color=s.FG_MUTED,
    )
    gap = Inches(0.3)
    win_w = (s.CONTENT_W - 2 * gap) / 3
    win_h = Inches(3.6)
    base_y = s.CONTENT_Y + Inches(0.65)
    screens = [
        ("Главная",            "sonar.app/home",       "ui-home.png"),
        ("Комната просмотра",  "sonar.app/room/r-42",  "ui-room.png"),
        ("Профиль",            "sonar.app/u/host",     "ui-profile.png"),
    ]
    img_dir = ROOT / "img"
    for i, (caption, url, fname) in enumerate(screens):
        x = s.PAD_X + i * (win_w + gap)
        # outer border
        add_rect(slide, x, base_y, win_w, win_h, s.BG_BASE)
        add_rect(slide, x, base_y, win_w, Emu(9525), s.BORDER_STRONG)
        add_rect(slide, x, base_y + win_h - Emu(9525), win_w, Emu(9525), s.BORDER_STRONG)
        add_rect(slide, x, base_y, Emu(9525), win_h, s.BORDER_STRONG)
        add_rect(slide, x + win_w - Emu(9525), base_y, Emu(9525), win_h, s.BORDER_STRONG)
        # chrome bar
        chrome_h = Inches(0.32)
        add_rect(slide, x, base_y, win_w, chrome_h, s.BG_MUTED)
        for j, color in enumerate([RGBColor(0xEF,0x44,0x44), RGBColor(0xF5,0x9E,0x0B), RGBColor(0x10,0xB9,0x81)]):
            dot_x = x + Inches(0.12) + j * Inches(0.18)
            dot_y = base_y + Inches(0.1)
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, dot_x, dot_y, Inches(0.13), Inches(0.13))
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.fill.background()
            dot.shadow.inherit = False
        # url
        add_textbox(
            slide, x + Inches(0.9), base_y + Inches(0.04), win_w - Inches(1.0), chrome_h - Inches(0.04),
            url,
            font=s.FONT_MONO, size=Pt(9), color=s.FG_MUTED, anchor=MSO_ANCHOR.MIDDLE,
        )
        # screenshot area
        inner_x = x + Emu(9525)
        inner_y = base_y + chrome_h
        inner_w = win_w - 2 * Emu(9525)
        inner_h = win_h - chrome_h - Emu(9525)
        png_path = img_dir / fname
        if png_path.exists():
            from PIL import Image as _PILImage
            with _PILImage.open(png_path) as im:
                iw, ih = im.size
            ratio = min(inner_w / (iw * 9525), inner_h / (ih * 9525))
            w = int(iw * 9525 * ratio)
            h = int(ih * 9525 * ratio)
            ix = inner_x + (inner_w - w) // 2
            iy = inner_y + (inner_h - h) // 2
            slide.shapes.add_picture(str(png_path), ix, iy, width=w, height=h)
        else:
            add_rect(slide, inner_x, inner_y, inner_w, inner_h, s.BG_ELEVATED)
            add_textbox(
                slide, inner_x, inner_y, inner_w, inner_h,
                f"[ скриншот: {caption.lower()} ]",
                font=s.FONT_MONO, size=Pt(10), color=s.FG_MUTED,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
            )
        # caption
        add_textbox(
            slide, x, base_y + win_h + Inches(0.1), win_w, Inches(0.35),
            caption,
            font=s.FONT_MONO, size=Pt(12), color=s.FG_MUTED,
            align=PP_ALIGN.CENTER,
        )


CODE_SAMPLE = """async def receive_json(self, content, **kwargs):
    \"\"\"Обработка play / pause / seek от ведущего.\"\"\"
    if not self.is_host:
        return
    action = content.get("action")
    position = float(content.get("position", 0.0))
    room = await self.update_room_state(
        position=position,
        is_playing=action in ("play", "resume"),
        action=action,
    )
    # broadcast подписчикам через Redis
    await self.channel_layer.group_send(self.group_name, {
        "type": "room.state",
        "position": room.playback_position,
        "is_playing": room.is_playing,
        "server_time": timezone.now().timestamp(),
    })"""


def slide_code(prs, page_num):
    slide = slide_inner(prs, page_num, "Пример реализации синхронизации")
    add_textbox(
        slide, s.PAD_X, s.CONTENT_Y, s.CONTENT_W, Inches(0.4),
        "Фрагмент WebSocket-консьюмера: обработка действий ведущего и рассылка состояния подписчикам.",
        font=s.FONT_SANS, size=s.BODY_SMALL, color=s.FG_MUTED,
    )
    # dark code block
    box_y = s.CONTENT_Y + Inches(0.65)
    box_h = Inches(4.5)
    add_rect(slide, s.PAD_X, box_y, s.CONTENT_W, box_h, s.FG_PRIMARY)
    # filename strip
    add_rect(slide, s.PAD_X, box_y, s.CONTENT_W, Inches(0.32), RGBColor(0x1F, 0x29, 0x37))
    add_textbox(
        slide, s.PAD_X + Inches(0.3), box_y, s.CONTENT_W - Inches(0.6), Inches(0.32),
        "rooms/consumers.py",
        font=s.FONT_MONO, size=Pt(11), color=RGBColor(0x9C, 0xA3, 0xAF),
        anchor=MSO_ANCHOR.MIDDLE,
    )
    # code body
    code_x = s.PAD_X + Inches(0.5)
    code_y = box_y + Inches(0.45)
    code_w = s.CONTENT_W - Inches(1.0)
    add_textbox(
        slide, code_x, code_y, code_w, box_h - Inches(0.55),
        CODE_SAMPLE,
        font=s.FONT_MONO, size=Pt(12), color=RGBColor(0xE5, 0xE7, 0xEB),
    )


def slide_testing(prs, page_num):
    slide = slide_inner(prs, page_num, "Тестирование")
    blocks = [
        ("Бэкенд: pytest + pytest-asyncio",
         "Юнит-тесты сервисов аутентификации, сериализаторов, парсеров каталога.\n"
         "Интеграционные тесты ASGI WebSocket-консьюмеров комнат."),
        ("Фронтенд: Vitest + Testing Library",
         "Юнит-тесты утилит синхронизации, форматирования времени, обработки ошибок API.\n"
         "Компонент-тесты ключевых элементов интерфейса (Button, MessageCard, TextField)."),
        ("Граничные случаи",
         "Рассинхрон ≥ 2 с включает жёсткую коррекцию плеера.\n"
         "Истёкший JWT обновляется по refresh-токену без выхода пользователя."),
    ]
    y = s.CONTENT_Y
    for label, body in blocks:
        add_textbox(
            slide, s.PAD_X, y, s.CONTENT_W, Inches(0.35),
            label,
            font=s.FONT_MONO, size=Pt(12), color=s.ACCENT, bold=True,
        )
        add_textbox(
            slide, s.PAD_X, y + Inches(0.35), s.CONTENT_W, Inches(1.0),
            body,
            font=s.FONT_SANS, size=Pt(15), color=s.FG_PRIMARY,
        )
        y += Inches(1.55)


def slide_deployment(prs, page_num):
    slide = slide_inner(prs, page_num, "Интеграция и развёртывание")
    items = [
        ("Контейнеризация",  "Multi-stage Dockerfile для бэкенда (Daphne ASGI) и фронтенда (nginx)."),
        ("Оркестрация",      "docker-compose с сервисами postgres, redis, backend, frontend."),
        ("Reverse proxy",    "Caddy с автоматическим выпуском TLS-сертификатов через Let's Encrypt."),
        ("Развёртывание",    "Идемпотентный скрипт deploy.sh: обновление на боевом сервере одной командой."),
        ("Конфигурация",     "Переменные окружения (.env), отдельные профили для dev и prod."),
        ("Стенд",            "Развёрнут на собственном сервере, доступен по доменному имени с HTTPS."),
    ]
    y = s.CONTENT_Y
    for label, val in items:
        add_textbox(
            slide, s.PAD_X, y, Inches(3.2), Inches(0.55),
            label,
            font=s.FONT_MONO, size=Pt(12), color=s.ACCENT, bold=True,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_textbox(
            slide, s.PAD_X + Inches(3.4), y, s.CONTENT_W - Inches(3.4), Inches(0.55),
            val,
            font=s.FONT_SANS, size=Pt(15), color=s.FG_PRIMARY,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        y += Inches(0.7)


def slide_conclusion(prs, page_num):
    slide = slide_inner(prs, page_num, "Заключение")
    # goal achieved block
    add_textbox(
        slide, s.PAD_X, s.CONTENT_Y, s.CONTENT_W, Inches(0.35),
        "Достижение цели",
        font=s.FONT_MONO, size=Pt(12), color=s.ACCENT, bold=True,
    )
    add_textbox(
        slide, s.PAD_X, s.CONTENT_Y + Inches(0.35), s.CONTENT_W, Inches(0.9),
        "Платформа потокового видео с синхронным совместным просмотром реализована,\nразвёрнута на боевом сервере, готова к демонстрации.",
        font=s.FONT_SANS, size=Pt(16), color=s.FG_PRIMARY,
    )
    # tasks checklist
    y2 = s.CONTENT_Y + Inches(1.7)
    add_textbox(
        slide, s.PAD_X, y2, s.CONTENT_W, Inches(0.35),
        "Выполнение задач",
        font=s.FONT_MONO, size=Pt(12), color=s.ACCENT, bold=True,
    )
    items = [
        "1.  Архитектура клиент-серверной системы спроектирована и реализована.",
        "2.  Потоковая отдача видео по HLS работает в современных браузерах.",
        "3.  WebSocket-синхронизация плеера работает с компенсацией дрейфа.",
        "4.  UI разработан, адаптивен для десктопа и мобильной версии.",
        "5.  Контейнерное развёртывание автоматизировано, HTTPS через Let's Encrypt.",
    ]
    add_textbox(
        slide, s.PAD_X, y2 + Inches(0.35), s.CONTENT_W, Inches(3.0),
        "\n".join(items),
        font=s.FONT_SANS, size=Pt(15), color=s.FG_PRIMARY,
    )


def slide_qa(prs):
    """Final 'questions' slide — same composition as title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_textbox(
        slide,
        s.PAD_X, s.HEADER_Y, Inches(6), Inches(0.3),
        "Sonar",
        font=s.FONT_MONO, size=Pt(13), color=s.ACCENT, bold=True,
    )
    add_textbox(
        slide,
        s.PAD_X, Inches(2.6), s.CONTENT_W, Inches(1.5),
        "Вопросы и ответы",
        font=s.FONT_SANS, size=Pt(54), bold=True, color=s.FG_PRIMARY,
    )
    add_rect(slide, s.PAD_X, Inches(4.0), Inches(1.0), Inches(0.06), s.ACCENT)
    # authors small
    y = Inches(4.3)
    for name, role in AUTHORS:
        add_textbox(
            slide, s.PAD_X, y, s.CONTENT_W, Inches(0.3),
            f"{name} · {role}",
            font=s.FONT_SANS, size=Pt(15), color=s.FG_MUTED,
        )
        y += Inches(0.4)
    # theme footer
    add_textbox(
        slide, s.PAD_X, s.SLIDE_H - Inches(0.6), s.CONTENT_W, Inches(0.3),
        THEME,
        font=s.FONT_MONO, size=Pt(11), color=s.FG_MUTED,
    )


# ---------- main ----------

def build():
    prs = Presentation()
    prs.slide_width = s.SLIDE_W
    prs.slide_height = s.SLIDE_H

    # 1 title
    slide_title(prs)
    # 2 goal & tasks
    slide_goal_tasks(prs, 2)
    # 3 roles
    slide_roles(prs, 3)
    # 4 means
    slide_means(prs, 4)
    img = ROOT / "img"
    # 5 gantt
    slide_image(prs, 5, "Планирование работ", png_path=img / "gantt.png")
    # 6 architecture
    slide_image(
        prs, 6, "Архитектура системы",
        "Контейнерная диаграмма в нотации C4: клиент, шлюз, сервисы, хранилища.",
        png_path=img / "c4-container.png",
    )
    # 7 ERD
    slide_image(prs, 7, "Диаграммы кода", png_path=img / "erd.png")
    # 8 UI tree
    slide_image(prs, 8, "Диаграмма компонентов пользовательского интерфейса", png_path=img / "ui-tree.png")
    # 9-11 парные слайды UI: Figma макет + реализация
    slide_ui_home(prs, 9)
    slide_ui_room(prs, 10)
    slide_ui_profile(prs, 11)
    # 12 sync algorithm
    slide_image(prs, 12, "Совместный просмотр: синхронизация плеера", png_path=img / "sync-sequence.png")
    # 13 code example
    slide_code(prs, 13)
    # 14 testing
    slide_testing(prs, 14)
    # 15 deployment
    slide_deployment(prs, 15)
    # 16 conclusion
    slide_conclusion(prs, 16)
    # 17 Q&A
    slide_qa(prs)

    prs.save(str(OUT_PPTX))
    print(f"saved: {OUT_PPTX}")
    print(f"slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()

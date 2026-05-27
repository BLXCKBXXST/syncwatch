"""
Первая итерация черновика презентации blxck.hub для ПИШ Иннополис.
- удаляет нерелевантные слайды (инструкция, дубль, ML-слайды)
- заполняет слайды, не требующие диаграмм/скриншотов

Запуск: /tmp/pptx-venv/bin/python fill_draft.py
"""
from copy import deepcopy
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

ROOT = Path(__file__).parent
SRC = ROOT / "template.pptx"
DST = ROOT / "draft.pptx"

BLACK = RGBColor(0x00, 0x00, 0x00)

TITLE_THEME = "Создание функциональной веб-платформы для потокового видео"
AUTHORS = [
    "Андреев Андрей Валерьевич, ст. преподаватель «СибГУТИ»",
    "Лукинов Виталий Леонидович, доцент «СибГУТИ»",
]


def delete_slide(prs, idx_zero_based):
    """Remove slide by zero-based index."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[idx_zero_based])


def set_text(text_frame, lines, font_size=None, color=BLACK):
    """Replace all text in a text frame with the given lines (preserves run formatting where possible)."""
    text_frame.clear()
    for i, line in enumerate(lines):
        para = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        run = para.add_run()
        run.text = line
        if font_size:
            run.font.size = Pt(font_size)
        run.font.color.rgb = color


def find_shape(slide, predicate):
    for sh in slide.shapes:
        if predicate(sh):
            return sh
    return None


def fill_title_slide(slide, theme, authors):
    """Slide layout 'Титульный слайд': big theme + up to 4 authors."""
    # Заголовок-тема: первый TEXT_BOX с заголовком
    theme_shape = find_shape(slide, lambda s: s.has_text_frame and s.text_frame.text.strip() == "Тема итоговой работы")
    if theme_shape is None:
        # fallback: any shape named "Заголовок 1"
        theme_shape = find_shape(slide, lambda s: s.name.startswith("Заголовок") and s.has_text_frame)
    if theme_shape is not None:
        set_text(theme_shape.text_frame, [theme])

    # Авторы: AUTO_SHAPE "Прямоугольник 5"
    authors_shape = find_shape(slide, lambda s: s.has_text_frame and "Фамилия Имя Отчество" in s.text_frame.text)
    if authors_shape is not None:
        # У шаблона 4 строки; заполняем имеющимися авторами, остальные оставляем пустыми
        lines = list(authors) + [""] * max(0, 4 - len(authors))
        # Сохраняем исходный размер шрифта первого run
        first_para = authors_shape.text_frame.paragraphs[0]
        original_size = None
        if first_para.runs:
            original_size = first_para.runs[0].font.size
        authors_shape.text_frame.clear()
        for i, line in enumerate(lines[:4]):
            para = authors_shape.text_frame.paragraphs[0] if i == 0 else authors_shape.text_frame.add_paragraph()
            run = para.add_run()
            run.text = line
            if original_size:
                run.font.size = original_size


def fill_goal_tasks(slide):
    """Slide 3: 'Цель и задачи проекта' — TextBox 6 has template lines."""
    box = find_shape(slide, lambda s: s.has_text_frame and "Цель проекта" in s.text_frame.text)
    if box is None:
        return
    lines = [
        "Цель проекта:",
        "Разработать функциональную веб-платформу для потокового видео",
        "с возможностью синхронного совместного просмотра.",
        "",
        "Задачи проекта:",
        "1. Спроектировать архитектуру системы: бэкенд, фронтенд, инфраструктура.",
        "2. Реализовать потоковую отдачу видео по протоколу HLS.",
        "3. Реализовать механизм синхронных комнат на WebSocket с серверной синхронизацией плеера.",
        "4. Разработать адаптивный пользовательский интерфейс (десктоп и мобайл).",
        "5. Обеспечить контейнерное развёртывание с автоматическим HTTPS.",
    ]
    set_text(box.text_frame, lines)


def fill_roles_table(slide):
    """Slide 4: таблица ФИО × Роль. В шаблоне 5 строк (1 заголовок + 4 ролей)."""
    table_shape = find_shape(slide, lambda s: s.has_table)
    if table_shape is None:
        return
    tbl = table_shape.table
    # Строки 1..4 — четыре роли. Колонка 0 — ФИО.
    assignments = [
        "Лукинов В.Л.",   # Моделирование архитектуры ИС
        "Лукинов В.Л.",   # Разработка бизнес-логики ИС
        "Андреев А.В.",   # Прототипирование и разработка UI
        "Андреев А.В.",   # Тестирование и интеграционные работы
    ]
    for row_idx, name in enumerate(assignments, start=1):
        cell = tbl.rows[row_idx].cells[0]
        cell.text_frame.clear()
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = name
        r.font.color.rgb = BLACK
        # Колонка 1 (роль) — установим чёрный цвет, в шаблоне текст был красный
        role_cell = tbl.rows[row_idx].cells[1]
        for para in role_cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.color.rgb = BLACK


def fill_content_placeholder(slide, lines):
    """Заполнить Content Placeholder (обычно стоит в центре с красным текстом-подсказкой)."""
    ph = find_shape(slide, lambda s: s.has_text_frame and (
        s.name.startswith("Content Placeholder") or s.name.startswith("Объект")
    ))
    if ph is None:
        return
    set_text(ph.text_frame, lines, color=BLACK)


def fill_means_of_implementation(slide):
    lines = [
        "Бэкенд: Python 3.10+, Django 5.1, Django REST Framework 3.15, Django Channels 4.1 (ASGI + WebSocket)",
        "База данных: PostgreSQL 16 (продакшен), SQLite (разработка)",
        "Кэш и pub/sub: Redis 7 (channel layer для WebSocket)",
        "Фронтенд: React 19.2, React Router 7, Vite 8, axios",
        "Потоковое видео: HLS (HTTP Live Streaming), hls.js, HTML5 video",
        "Инфраструктура: Docker, docker-compose, Caddy (reverse proxy и auto-HTTPS Let's Encrypt)",
        "Тестирование: pytest + pytest-asyncio (бэк), Vitest + Testing Library (фронт)",
    ]
    fill_content_placeholder(slide, lines)


def fill_testing(slide):
    lines = [
        "Бэкенд (pytest + pytest-asyncio):",
        "  – юнит-тесты сервисов аутентификации, сериализаторов, парсеров каталога;",
        "  – интеграционные тесты ASGI WebSocket-консьюмеров комнат.",
        "",
        "Фронтенд (Vitest + @testing-library/react):",
        "  – юнит-тесты утилит синхронизации (sync.js), форматирования времени, обработки ошибок API;",
        "  – компонент-тесты ключевых элементов UI (Button, MessageCard, TextField).",
        "",
        "Граничные значения: рассинхрон ≥ 2 с включает жёсткую коррекцию плеера; истёкший JWT обновляется по refresh-токену.",
    ]
    fill_content_placeholder(slide, lines)


def fill_deployment(slide):
    lines = [
        "Контейнеризация: multi-stage Dockerfile для бэкенда (Daphne ASGI) и фронтенда (nginx).",
        "Оркестрация: docker-compose с сервисами postgres, redis, backend, frontend.",
        "Reverse proxy: Caddy с автоматическим выпуском TLS-сертификатов через Let's Encrypt.",
        "Развёртывание: скрипт deploy.sh — идемпотентное обновление на боевом сервере.",
        "Конфигурация: переменные окружения (.env), отдельные профили для dev и prod.",
        "Стенд: развёрнут на собственном сервере, доступен по доменному имени с HTTPS.",
    ]
    fill_content_placeholder(slide, lines)


def fill_conclusion(slide):
    """Slide 18: 'Заключение' — TextBox 10 has cell-template."""
    box = find_shape(slide, lambda s: s.has_text_frame and "Достижение цели" in s.text_frame.text)
    if box is None:
        return
    lines = [
        "Достижение цели:",
        "Платформа потокового видео с синхронным совместным просмотром реализована,",
        "развёрнута на боевом сервере, готова к демонстрации.",
        "",
        "Выполнение задач:",
        "1. Архитектура клиент-серверной системы спроектирована и реализована.",
        "2. Потоковая отдача видео по HLS работает в современных браузерах.",
        "3. WebSocket-синхронизация плеера работает с компенсацией дрейфа.",
        "4. UI разработан, адаптивен для десктопа и мобильной версии.",
        "5. Контейнерное развёртывание автоматизировано, HTTPS — Let's Encrypt.",
    ]
    set_text(box.text_frame, lines)


# --- main ---

prs = Presentation(str(SRC))
print(f"loaded: {len(prs.slides)} slides")

# Удаляем в обратном порядке (0-индексация): 13 (бывший 14, ML), 12 (бывший 13, ML), 10 (бывший 11, дубль), 1 (бывший 2, инструкция)
for i in [13, 12, 10, 1]:
    delete_slide(prs, i)
print(f"after deletions: {len(prs.slides)} slides")

# Новая нумерация (после удалений) 0-индексация:
# 0: Титул
# 1: Цель и задачи
# 2: Распределение ролей
# 3: Средства реализации
# 4: Планирование работ            -> позже (Гант)
# 5: Архитектура системы           -> позже (C4 diagram)
# 6: Диаграммы кода                -> позже (UML+ERD)
# 7: Диаграмма компонентов UI      -> позже
# 8: Прототипирование UI           -> позже (figma скрины)
# 9: Алгоритм валидации данных     -> переименовать в "Алгоритм синхронизации плеера" + диаграмма
# 10: Пример функции Счётчик       -> позже (код-скрин)
# 11: Тестирование
# 12: Интеграция и развертывание
# 13: Заключение
# 14: Вопросы-ответы (повтор титула)

fill_title_slide(prs.slides[0], TITLE_THEME, AUTHORS)
fill_goal_tasks(prs.slides[1])
fill_roles_table(prs.slides[2])
fill_means_of_implementation(prs.slides[3])
fill_testing(prs.slides[11])
fill_deployment(prs.slides[12])
fill_conclusion(prs.slides[13])
fill_title_slide(prs.slides[14], TITLE_THEME, AUTHORS)

prs.save(str(DST))
print(f"saved: {DST}")
